#!/usr/bin/env python3
"""Export Instagram 9:16 drill slides to PNG, MP4, and PPTX."""

from __future__ import annotations

import subprocess
import sys
import time
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
HTML = ROOT / "docs" / "instagram-drill-reel.html"
OUT = Path(__file__).resolve().parent
SLIDES = OUT / "slides"
CHROME = "/usr/bin/google-chrome"
SLIDE_COUNT = 8
HOLD_SEC = 2.5


def run(cmd: list[str]) -> None:
    print("+", " ".join(cmd), flush=True)
    subprocess.run(cmd, check=True)


def export_pngs() -> list[Path]:
    SLIDES.mkdir(parents=True, exist_ok=True)
    html_uri = HTML.resolve().as_uri()
    paths: list[Path] = []
    for n in range(1, SLIDE_COUNT + 1):
        dest = SLIDES / f"{n:02d}.png"
        run(
            [
                CHROME,
                "--headless=new",
                "--disable-gpu",
                "--hide-scrollbars",
                "--no-sandbox",
                "--force-device-scale-factor=1",
                "--window-size=1080,1920",
                f"--virtual-time-budget=8000",
                f"--screenshot={dest}",
                f"{html_uri}?export={n}",
            ]
        )
        if not dest.exists() or dest.stat().st_size < 10_000:
            raise SystemExit(f"Screenshot failed: {dest}")
        paths.append(dest)
        time.sleep(0.2)
    return paths


def export_mp4(pngs: list[Path]) -> Path:
    dest = OUT / "mentor-maths-daily-drill-reel.mp4"
    # Hold each slide HOLD_SEC with a short crossfade.
    fade = 0.25
    inputs: list[str] = []
    for png in pngs:
        inputs.extend(["-loop", "1", "-t", f"{HOLD_SEC + fade}", "-i", str(png)])

    filters = []
    last = "0:v"
    # Scale/format each input, then chain xfades.
    for i in range(len(pngs)):
        filters.append(f"[{i}:v]scale=1080:1920:force_original_aspect_ratio=disable,format=yuv420p,setsar=1[v{i}]")
    last = "v0"
    offset = HOLD_SEC
    for i in range(1, len(pngs)):
        out = f"x{i}"
        filters.append(f"[{last}][v{i}]xfade=transition=fade:duration={fade}:offset={offset:.2f}[{out}]")
        last = out
        offset += HOLD_SEC

    run(
        [
            "ffmpeg",
            "-y",
            *inputs,
            "-filter_complex",
            ";".join(filters),
            "-map",
            f"[{last}]",
            "-c:v",
            "libx264",
            "-pix_fmt",
            "yuv420p",
            "-r",
            "30",
            "-movflags",
            "+faststart",
            str(dest),
        ]
    )
    return dest


def export_pptx(pngs: list[Path]) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        run([sys.executable, "-m", "pip", "install", "--quiet", "python-pptx"])
        from pptx import Presentation
        from pptx.util import Inches

    dest = OUT / "mentor-maths-daily-drill.pptx"
    prs = Presentation()
    # Portrait 1080x1920 at 96 dpi → 11.25" × 20"
    prs.slide_width = Inches(11.25)
    prs.slide_height = Inches(20)
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    prs.save(dest)
    return dest


def main() -> None:
    if not Path(CHROME).exists():
        raise SystemExit(f"Chrome not found at {CHROME}")
    pngs = export_pngs()
    mp4 = export_mp4(pngs)
    pptx = export_pptx(pngs)
    print("Wrote:")
    for p in [*pngs, mp4, pptx]:
        print(f"  {p} ({p.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
