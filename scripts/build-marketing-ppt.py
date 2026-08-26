"""Build MentorMaths coaching-class marketing PowerPoint."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "MentorMaths_Coaching_Marketing.pptx"
ASSETS = Path(
    r"C:\Users\Atul.Gupta\.cursor\projects\c-Users-Atul-Gupta-maths-foundation\assets"
)

SHOTS = {
    "class": ASSETS
    / "c__Users_Atul.Gupta_AppData_Roaming_Cursor_User_workspaceStorage_ad0ab84d653cd53839d2ffba3a60bed5_images_image-0aad162b-3c07-417a-8bca-2234a80d30ff.png",
    "study": ASSETS
    / "c__Users_Atul.Gupta_AppData_Roaming_Cursor_User_workspaceStorage_ad0ab84d653cd53839d2ffba3a60bed5_images_image-be5f5df5-191d-4ea2-aefd-43becf1a111c.png",
    "coverage": ASSETS
    / "c__Users_Atul.Gupta_AppData_Roaming_Cursor_User_workspaceStorage_ad0ab84d653cd53839d2ffba3a60bed5_images_image-cbf36cc1-44f0-495a-b708-3b204596f770.png",
    "basics": ASSETS
    / "c__Users_Atul.Gupta_AppData_Roaming_Cursor_User_workspaceStorage_ad0ab84d653cd53839d2ffba3a60bed5_images_image-7e37a4f7-7aef-45b5-904f-9c2845b68c0f.png",
    "email": ASSETS
    / "c__Users_Atul.Gupta_AppData_Roaming_Cursor_User_workspaceStorage_ad0ab84d653cd53839d2ffba3a60bed5_images_image-9c726a82-f7eb-4e3f-98fc-e7eadb8d77a0.png",
}

# Brand (matches product: deep indigo + teal, not generic purple gradient)
NAVY = RGBColor(0x1E, 0x1B, 0x4B)
INK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x4B, 0x55, 0x63)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
CORAL = RGBColor(0xC2, 0x41, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
GREEN = RGBColor(0x15, 0x80, 0x3D)


def set_run(run, text, size=18, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color, font=font)
    return box


def fill_shape(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_shape(shape, color)
    return shape


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def header_bar(slide, title, subtitle=None):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.95), NAVY)
    add_textbox(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.4), title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.35), subtitle, size=13, color=RGBColor(0xC7, 0xD2, 0xFE))


def footer(slide, page, total=12):
    add_textbox(
        slide,
        Inches(0.5),
        Inches(7.15),
        Inches(10),
        Inches(0.3),
        "mentormaths.in  ·  Plan · Practice · Perform",
        size=11,
        color=MUTED,
    )
    add_textbox(
        slide,
        Inches(11.5),
        Inches(7.15),
        Inches(1.3),
        Inches(0.3),
        f"{page}/{total}",
        size=11,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_bullets(slide, left, top, width, height, items, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        run = p.add_run()
        set_run(run, f"•  {item}", size=size, color=INK)
    return box


def fit_picture(slide, path: Path, left, top, max_w, max_h):
    if not path.exists():
        add_textbox(slide, left, top, max_w, Inches(0.4), f"[Missing: {path.name}]", size=12, color=CORAL)
        return None
    pic = slide.shapes.add_picture(str(path), left, top)
    # scale to fit
    w, h = pic.width, pic.height
    scale = min(max_w / w, max_h / h)
    pic.width = int(w * scale)
    pic.height = int(h * scale)
    return pic


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = 12

    # 1 Title
    s = blank_slide(prs)
    rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    add_textbox(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5), "MENTORMATHS", size=18, bold=True, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(1), "Plan · Practice · Perform", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(
        s,
        Inches(1.5),
        Inches(3.5),
        Inches(10.3),
        Inches(1),
        "The maths coaching OS — exam plans, daily drills, revision queues,\nand parent alerts built by a teacher who runs real classes.",
        size=18,
        color=RGBColor(0xE0, 0xE7, 0xFF),
        align=PP_ALIGN.CENTER,
    )
    add_textbox(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4), "For coaching classes & serious home mentors  ·  mentormaths.in", size=14, color=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.CENTER)

    # 2 Punch lines
    s = blank_slide(prs)
    header_bar(s, "One line coaches remember", "Punch lines you can put on a flyer or WhatsApp broadcast")
    punches = [
        ("Not another question bank.", "A full coaching operating system — plan → practice → exam → revise."),
        ("Wrong answers never die quietly.", "Every miss enters a revision queue until the student redoes it right."),
        ("Parents see progress without chasing you.", "Email + WhatsApp when work is done — score, time, and what went wrong."),
        ("Built in the classroom, not a boardroom.", "Made by a practising mentor · shipped weekly with agile changes for real batches."),
        ("TrainingPeaks for Maths.", "Syllabus plan + daily workout + race day (school exam) + performance stats."),
    ]
    y = 1.2
    for title, body in punches:
        rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(0.95), LIGHT)
        add_textbox(s, Inches(0.7), Inches(y + 0.12), Inches(11.8), Inches(0.35), title, size=18, bold=True, color=TEAL)
        add_textbox(s, Inches(0.7), Inches(y + 0.48), Inches(11.8), Inches(0.4), body, size=14, color=INK)
        y += 1.05
    footer(s, 2, total)

    # 3 Problem
    s = blank_slide(prs)
    header_bar(s, "The coaching-class problem", "Generic apps don’t run a batch the way you do")
    left_items = [
        "WhatsApp groups full of PDFs — no completion truth",
        "Excel trackers go stale every exam week",
        "Students “finish” but wrong sums never get redone",
        "Parents ask “how is he doing?” — you dig for answers",
        "One syllabus, many exam dates per child",
    ]
    right_items = [
        "Live class dashboard: completion, score, hours, revision",
        "Per-student exam plan tied to chapters",
        "Revision queue: pending + wrong until cleared",
        "Auto email / WhatsApp on every submission",
        "Content bank: practice, test, written, FIB, formula, books",
    ]
    add_textbox(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.4), "Today’s pain", size=20, bold=True, color=CORAL)
    add_bullets(s, Inches(0.5), Inches(1.7), Inches(5.8), Inches(4.5), left_items, size=15)
    add_textbox(s, Inches(7), Inches(1.2), Inches(5.8), Inches(0.4), "With MentorMaths", size=20, bold=True, color=TEAL)
    add_bullets(s, Inches(7), Inches(1.7), Inches(5.8), Inches(4.5), right_items, size=15)
    footer(s, 3, total)

    # 4 Coaching dashboard screenshot
    s = blank_slide(prs)
    header_bar(s, "Coaching-class dashboard", "Class status & exam plans — who is ready for the next paper")
    fit_picture(s, SHOTS["class"], Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.5))
    footer(s, 4, total)

    # 5 What coaches see
    s = blank_slide(prs)
    header_bar(s, "What the coach sees at a glance", "From the class screen — actionable, not decorative")
    cards = [
        ("Completion %", "Sets done vs planned for studied chapters"),
        ("Score %", "Average across scored attempts"),
        ("Revision", "Pending + wrong still to redo"),
        ("Days & hours", "Login streak and active study time"),
        ("Exam plan", "Half-yearly / unit test date per child"),
        ("Alerts", "“4 students have no upcoming exam plan”"),
    ]
    positions = [(0.5, 1.3), (4.5, 1.3), (8.5, 1.3), (0.5, 3.6), (4.5, 3.6), (8.5, 3.6)]
    for (title, body), (x, y) in zip(cards, positions):
        rect(s, Inches(x), Inches(y), Inches(3.7), Inches(1.9), LIGHT)
        add_textbox(s, Inches(x + 0.2), Inches(y + 0.25), Inches(3.3), Inches(0.4), title, size=18, bold=True, color=NAVY)
        add_textbox(s, Inches(x + 0.2), Inches(y + 0.8), Inches(3.3), Inches(0.9), body, size=14, color=MUTED)
    footer(s, 5, total)

    # 6 Student study plan
    s = blank_slide(prs)
    header_bar(s, "Student view — study plan performance", "Same truth the coach sees · chapters tagged to exam date")
    fit_picture(s, SHOTS["study"], Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.5))
    footer(s, 6, total)

    # 7 Drill concept
    s = blank_slide(prs)
    header_bar(s, "The drill concept", "Mastery = finish the set + clear every wrong sum")
    add_bullets(
        s,
        Inches(0.5),
        Inches(1.25),
        Inches(6.2),
        Inches(5),
        [
            "Guided practice: one sum at a time, hints, give-up tracked",
            "Learner / Achiever / Expert mixes — not one-size worksheets",
            "Wrong answers → revision queue (pending until redone)",
            "Basics drill: tables / squares / cubes with seconds-per-blank",
            "Formula drill: daily spaced practice from your bank",
            "Chapter tests autosave — WiFi drop doesn’t wipe progress",
        ],
        size=16,
    )
    fit_picture(s, SHOTS["basics"], Inches(6.9), Inches(1.25), Inches(5.9), Inches(5.2))
    footer(s, 7, total)

    # 8 Content coverage
    s = blank_slide(prs)
    header_bar(s, "Content coverage — your bank, chapter × format", "Practice · Test · Written · Fill-blank · Formula · Books")
    fit_picture(s, SHOTS["coverage"], Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.5))
    footer(s, 8, total)

    # 9 Notifications
    s = blank_slide(prs)
    header_bar(s, "Email & WhatsApp — parents stay in the loop", "Every submission: score, time, attempt history, what went wrong")
    fit_picture(s, SHOTS["email"], Inches(0.4), Inches(1.15), Inches(7.8), Inches(5.5))
    add_textbox(s, Inches(8.5), Inches(1.4), Inches(4.3), Inches(0.4), "In the mail", size=18, bold=True, color=TEAL)
    add_bullets(
        s,
        Inches(8.5),
        Inches(1.9),
        Inches(4.3),
        Inches(4.5),
        [
            "Attempt # and score",
            "Time taken vs target date",
            "All prior attempts compared",
            "First-try correct / hints used",
            "“What went wrong” review list",
            "WhatsApp templates for assignments too",
        ],
        size=14,
    )
    footer(s, 9, total)

    # 10 Competitive
    s = blank_slide(prs)
    header_bar(
        s,
        "How we differ from market tools",
        "Byju’s · EduGain · Open Door · fee CRMs — vs a coaching-class OS",
    )

    headers = ["Capability", "Byju’s", "EduGain", "Open Door", "Fee / tuition\nCRMs", "MentorMaths"]
    rows = [
        [
            "Exam-linked study plan",
            "Course packs",
            "Topic practice",
            "Topic tests",
            "Rare",
            "Yes — per student",
        ],
        [
            "Revision of wrong sums",
            "In-app only",
            "Focus areas",
            "Assessment gap",
            "—",
            "Auto redo queue",
        ],
        [
            "Basics / formula drills",
            "General app",
            "Worksheets",
            "—",
            "—",
            "Class-tuned daily",
        ],
        [
            "Parent email / WhatsApp",
            "Parent app",
            "Limited",
            "School reports",
            "Fees only",
            "On every attempt",
        ],
        [
            "Coach class dashboard",
            "Consumer-led",
            "Teacher assign",
            "School admin",
            "Attendance/fees",
            "Done · score · hours",
        ],
        [
            "Your books + sets matrix",
            "Their content",
            "Generated bank",
            "Their workbooks",
            "—",
            "Chapter × format",
        ],
        [
            "Built for YOUR batch",
            "Big roadmap",
            "Product SKU",
            "School program",
            "Vendor queue",
            "Weekly — mentor builds",
        ],
    ]

    # 6 columns across 13.333" slide
    col_w = [2.55, 1.95, 1.95, 1.95, 1.95, 2.35]
    col_x = [0.35]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w + 0.08)
    last_i = len(headers) - 1

    y = 1.15
    for i, h in enumerate(headers):
        if i == 0:
            fill = TEAL
        elif i == last_i:
            fill = RGBColor(0x14, 0x5A, 0x32)
        else:
            fill = NAVY
        rect(s, Inches(col_x[i]), Inches(y), Inches(col_w[i]), Inches(0.72), fill)
        add_textbox(
            s,
            Inches(col_x[i] + 0.05),
            Inches(y + 0.08),
            Inches(col_w[i] - 0.1),
            Inches(0.58),
            h,
            size=11,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

    y = 1.92
    for r, row in enumerate(rows):
        bg = LIGHT if r % 2 == 0 else WHITE
        for i, cell in enumerate(row):
            cell_bg = RGBColor(0xEC, 0xFE, 0xFF) if i == 0 else (RGBColor(0xEC, 0xFD, 0xF5) if i == last_i else bg)
            rect(s, Inches(col_x[i]), Inches(y), Inches(col_w[i]), Inches(0.58), cell_bg)
            color = GREEN if i == last_i else INK
            add_textbox(
                s,
                Inches(col_x[i] + 0.05),
                Inches(y + 0.12),
                Inches(col_w[i] - 0.1),
                Inches(0.4),
                cell,
                size=11,
                bold=(i == 0 or i == last_i),
                color=color,
                align=PP_ALIGN.CENTER if i else PP_ALIGN.LEFT,
            )
        y += 0.58

    add_textbox(
        s,
        Inches(0.4),
        Inches(6.75),
        Inches(12.5),
        Inches(0.3),
        "Fair positioning by typical product strength — not a full SKU audit of every plan.",
        size=10,
        color=MUTED,
    )
    footer(s, 10, total)

    # 11 Built by user + agile
    s = blank_slide(prs)
    header_bar(s, "Built by the mentor · shipped with agile speed", "Not a distant SaaS — tuned for the classes that use it")
    blocks = [
        ("Practitioner-built", "Designed while teaching real CBSE batches — features come from classroom friction, not feature lists."),
        ("Agile weekly releases", "Need a new drill range, exam label, or WhatsApp flow? Ship it for your class the same week."),
        ("Your syllabus, your books", "Ganita Prakash, RD Sharma, Greya Lakshmi… mapped chapter-wise into one coverage matrix."),
        ("Integrity + continuity", "Test protection, autosave on WiFi drops, Continue where the student left off."),
    ]
    y = 1.25
    for title, body in blocks:
        rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.15), LIGHT)
        add_textbox(s, Inches(0.75), Inches(y + 0.2), Inches(11.8), Inches(0.35), title, size=18, bold=True, color=NAVY)
        add_textbox(s, Inches(0.75), Inches(y + 0.55), Inches(11.8), Inches(0.45), body, size=14, color=MUTED)
        y += 1.3
    footer(s, 11, total)

    # 12 Close
    s = blank_slide(prs)
    rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    add_textbox(s, Inches(1), Inches(2.0), Inches(11.3), Inches(0.6), "Ready for your next batch?", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(
        s,
        Inches(1.5),
        Inches(2.9),
        Inches(10.3),
        Inches(1.2),
        "Run one class on MentorMaths for one exam cycle.\nWatch completion, revision, and parent calm improve together.",
        size=18,
        color=RGBColor(0xE0, 0xE7, 0xFF),
        align=PP_ALIGN.CENTER,
    )
    add_textbox(s, Inches(1), Inches(4.5), Inches(11.3), Inches(0.5), "mentormaths.in", size=28, bold=True, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.3), Inches(11.3), Inches(0.4), "Plan · Practice · Perform", size=16, color=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")
    for k, p in SHOTS.items():
        print(f"  {k}: {'OK' if p.exists() else 'MISSING'} {p.name}")


if __name__ == "__main__":
    build()
